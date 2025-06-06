import os
import csv
from pptx import Presentation
import sqlite3


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slide_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text.strip())
        slide_text.append("\n".join(texts))
    return slide_text


def build_database(pptx_dir, db_path="project_status.db"):
    conn = sqlite3.connect(db_path)
    cur = conn.cursor()
    cur.execute(
        "CREATE TABLE IF NOT EXISTS project_status ("\
        "id INTEGER PRIMARY KEY AUTOINCREMENT, "\
        "project_name TEXT, "\
        "slide_number INTEGER, "\
        "content TEXT)"
    )
    conn.commit()

    for fname in os.listdir(pptx_dir):
        if not fname.lower().endswith(".pptx"):
            continue
        pptx_path = os.path.join(pptx_dir, fname)
        slides = extract_text_from_pptx(pptx_path)
        project_name = os.path.splitext(fname)[0]
        for i, text in enumerate(slides, start=1):
            cur.execute(
                "INSERT INTO project_status (project_name, slide_number, content) "\
                "VALUES (?, ?, ?)",
                (project_name, i, text)
            )
    conn.commit()
    conn.close()


def export_csv(db_path="project_status.db", csv_path="project_status.csv"):
    conn = sqlite3.connect(db_path)
    cur = conn.cursor()
    cur.execute("SELECT project_name, slide_number, content FROM project_status ORDER BY project_name, slide_number")
    rows = cur.fetchall()
    conn.close()

    with open(csv_path, "w", newline="", encoding="utf-8") as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(["project_name", "slide_number", "content"])
        writer.writerows(rows)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Build a project status database from PPTX files.")
    parser.add_argument("--pptx_dir", required=True, help="Directory containing PPTX files")
    parser.add_argument("--db", default="project_status.db", help="Output SQLite DB path")
    parser.add_argument("--csv", help="Optional CSV export path")

    args = parser.parse_args()

    build_database(args.pptx_dir, args.db)

    if args.csv:
        export_csv(args.db, args.csv)
